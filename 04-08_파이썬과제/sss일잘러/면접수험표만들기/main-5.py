from pptx import Presentation 
from pptx.util import Pt
import pandas as pd
import copy
import os
os.chdir(os.path.dirname(os.path.abspath(__file__)))

df_from_excel = pd.read_excel("수험번호.xlsx")
prs = Presentation("수험표_샘플.pptx")

이름 = iter(df_from_excel['이름'])
수험번호 = iter(df_from_excel['수험번호'])
응시자수 = len(df_from_excel)

def copy_slide(prs, org_slide):
    # 새 슬라이드를 생성하고, 원본 슬라이드의 레이아웃을 사용합니다.
    copied_slide = prs.slides.add_slide(org_slide.slide_layout)
    
    # 원본 슬라이드의 모든 도형을 순회하며 복사합니다.
    for shape in org_slide.shapes:
        org_el = shape.element  # 원본 도형의 XML 요소를 가져옵니다.
        new_el = copy.deepcopy(org_el)  # 원본 도형의 깊은 복사본을 생성합니다.
        # 복사된 도형을 새 슬라이드에 추가합니다. 'p:extLst' 앞에 삽입하여 순서를 유지합니다.
        copied_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    
    # 원본 슬라이드의 모든 관계(relations)를 순회합니다.
    for rel in org_slide.part.rels.values():  # .values() 메서드를 사용해 실제 관계 객체들을 가져옵니다.
        # 노트 슬라이드와 관련된 관계는 제외합니다.
        if "notesSlide" not in rel.reltype:
            # 현재 관계의 타입(reltype)과 타겟(_target)을 사용하여 새 슬라이드에 동일한 관계를 추가하거나 가져옵니다.
            copied_slide.part.rels.get_or_add(
                rel.reltype,  # 관계의 타입 (예: 이미지, 차트 등에 대한 관계)
                rel._target   # 관계의 대상 (예: 이미지 파일, 차트 데이터 파일 등)
            )
    return copied_slide  # 수정된 새 슬라이드 객체를 반환합니다.


for i in range(int(응시자수/2)):
    next_slide = copy_slide(prs, prs.slides[0]) 
    for shape in next_slide.shapes:
        if shape.shape_type == 17 and shape.text == "이름":
            shape.text_frame.paragraphs[0].text = next(이름)
            shape.text_frame.paragraphs[0].font.size = Pt(40)
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].font.name = "나눔고딕"
        if shape.shape_type == 17 and shape.text == "수험번호":
            shape.text_frame.paragraphs[0].text = next(수험번호)
            shape.text_frame.paragraphs[0].font.size = Pt(40)
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].font.name = "나눔고딕"

prs.save("수험표_결과.pptx")
