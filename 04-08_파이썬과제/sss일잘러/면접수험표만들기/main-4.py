from pptx import Presentation 
from pptx.util import Pt
import os
os.chdir(os.path.dirname(os.path.abspath(__file__)))

prs = Presentation("수험표_샘플.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.shape_type == 17:
            print(shape.text)
